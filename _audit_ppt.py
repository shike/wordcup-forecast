"""Audit the generated PPT for layout problems."""
from pptx import Presentation
from pptx.util import Emu

prs = Presentation("输出/加拿大_对阵_波黑_2026-06-13.pptx")
SLIDE_W = 13.333
SLIDE_H = 7.5

print(f"Total pages: {len(prs.slides)}")
print(f"Slide W={SLIDE_W}, H={SLIDE_H}\n")

for i, slide in enumerate(prs.slides, 1):
    shapes_info = []
    for shape in slide.shapes:
        if shape.left is None or shape.top is None:
            continue
        l = shape.left.inches
        t = shape.top.inches
        w = shape.width.inches if shape.width else 0
        h = shape.height.inches if shape.height else 0
        right = l + w
        bottom = t + h
        text = ""
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()[:30]
        shapes_info.append((l, t, right, bottom, w, h, text, shape.shape_type))

    if not shapes_info:
        continue

    issues = []

    # 1. Off-page shapes
    for l, t, r, b, w, h, txt, _ in shapes_info:
        if r > SLIDE_W + 0.05 or b > SLIDE_H + 0.05 or l < -0.05 or t < -0.05:
            issues.append(f"OFF-PAGE: [{l:.1f},{t:.1f} {w:.1f}x{h:.1f}] '{txt}'")

    # 2. Overlapping text shapes (same z-order, both have text, overlap > 50%)
    text_shapes = [s for s in shapes_info if s[6]]
    for i in range(len(text_shapes)):
        for j in range(i + 1, len(text_shapes)):
            l1, t1, r1, b1, _, _, txt1, _ = text_shapes[i]
            l2, t2, r2, b2, _, _, txt2, _ = text_shapes[j]
            # bounding box overlap area
            ox = max(0, min(r1, r2) - max(l1, l2))
            oy = max(0, min(b1, b2) - max(t1, t2))
            overlap_area = ox * oy
            area1 = (r1 - l1) * (b1 - t1)
            area2 = (r2 - l2) * (b2 - t2)
            if area1 > 0 and area2 > 0:
                overlap_pct = overlap_area / min(area1, area2)
                if overlap_pct > 0.6 and len(txt1) > 2 and len(txt2) > 2:
                    issues.append(
                        f"OVERLAP {overlap_pct:.0%}: '{txt1}' vs '{txt2}'"
                    )

    if issues:
        print(f"\n=== Page {i} ({len(shapes_info)} shapes) — {len(issues)} issue(s) ===")
        for iss in issues[:15]:
            print(f"  • {iss}")
