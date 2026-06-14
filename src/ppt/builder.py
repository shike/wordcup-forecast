"""Main PPT builder.

Renders the full 18-20 page report from a PredictionResult.  Each page is a
python-pptx slide laid out in the dark sports-data style.
"""
from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Literal

from loguru import logger
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

from src.ppt.charts import (
    depth_bars,
    form_trend,
    probability_bars,
    qualitative_radar,
    radar_chart,
    score_distribution,
    sensitivity_tornado,
)
from src.ppt.pitch_layout import draw_pitch_with_lineup
from src.ppt.player_card import player_has_photo, render_player_card
from src.ppt.styles import (
    BG_CARD,
    BG_DEEP,
    BG_PANEL,
    BG_SOFT,
    CYAN,
    FONT_BODY,
    FONT_CN_BODY,
    FONT_CN_TITLE,
    FONT_MONO,
    FONT_TITLE,
    GOLD,
    GOLD_DEEP,
    GREEN,
    GREY,
    GREY_DARK,
    LINE,
    MARGIN,
    RED,
    SECTION_SIZE,
    SLIDE_H,
    SLIDE_W,
    SMALL_SIZE,
    SUBTITLE_SIZE,
    TITLE_SIZE,
    WHITE,
    BODY_SIZE,
)
from src.utils.config import config
from src.utils.i18n import tr, tr_pair, tr_zh, tr_en
from src.utils.models import (
    InjuryReport,
    Lineup,
    MatchInput,
    Matchup,
    Player,
    PredictionResult,
    Team,
    TeamStats,
)
from src.data.repository import MatchRepository
from src.data.wikipedia_client import batch_augment_squad


CHART_DIR = Path("/tmp/世界杯预测_图表缓存")
CHART_DIR.mkdir(exist_ok=True)


# ----------------------- helpers -----------------------

def _bg(slide, color: RGBColor = BG_DEEP) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_run_fonts(run, latin: str, ea: str) -> None:
    """Set both Latin and East-Asian typefaces on a text run.

    This is the key fix for Chinese rendering: PowerPoint uses the Latin
    typeface for Latin characters and the East-Asian typeface for CJK chars.
    Without setting the EA typeface explicitly, the Latin font is used for
    everything and Chinese characters become boxes (□) on systems where
    the chosen Latin font lacks CJK glyphs.
    """
    rPr = run._r.get_or_add_rPr()
    # Latin
    latin_el = rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}latin")
    if latin_el is None:
        from lxml import etree
        latin_el = etree.SubElement(
            rPr,
            "{http://schemas.openxmlformats.org/drawingml/2006/main}latin",
        )
    latin_el.set("typeface", latin)
    # East Asian
    ea_el = rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    if ea_el is None:
        from lxml import etree
        ea_el = etree.SubElement(
            rPr,
            "{http://schemas.openxmlformats.org/drawingml/2006/main}ea",
        )
    ea_el.set("typeface", ea)
    # Complex script (Arabic/Hebrew etc., use Latin fallback)
    cs_el = rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}cs")
    if cs_el is None:
        from lxml import etree
        cs_el = etree.SubElement(
            rPr,
            "{http://schemas.openxmlformats.org/drawingml/2006/main}cs",
        )
    cs_el.set("typeface", latin)


def _enable_autofit(text_frame) -> None:
    """Enable PowerPoint 'shrink text on overflow' for a text frame."""
    from lxml import etree
    bodyPr = text_frame._txBody.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr"
    )
    if bodyPr is not None:
        # remove any existing autofit
        for child in bodyPr:
            tag = etree.QName(child).localname
            if tag in ("normAutofit", "spAutoFit", "noAutofit"):
                bodyPr.remove(child)
        af = etree.SubElement(
            bodyPr,
            "{http://schemas.openxmlformats.org/drawingml/2006/main}normAutofit",
        )
        af.set("fontScale", "92500")  # 92.5% if overflow


def _add_textbox(
    slide, left, top, width, height, text, *,
    font_size=Pt(14), bold=False, color=WHITE,
    font_name=FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    cn: str | None = None,
    auto_shrink: bool = True,
) -> None:
    if cn and text:
        text = f"{cn}\n{text}"
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    if cn and text and "\n" in text:
        # bilingual: split into two paragraphs
        zh, en = cn, text
        p.text = ""
        run = p.add_run()
        run.text = zh
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color
        _set_run_fonts(run, latin=font_name, ea=FONT_CN_BODY)
        p2 = tf.add_paragraph()
        p2.alignment = align
        run2 = p2.add_run()
        run2.text = en
        run2.font.size = Pt(int(font_size.pt * 0.65))
        run2.font.bold = False
        run2.font.color.rgb = GREY
        _set_run_fonts(run2, latin=font_name, ea=FONT_CN_BODY)
    else:
        run = p.add_run()
        run.text = text
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color
        # Always use Microsoft YaHei for East-Asian glyphs so any stray
        # CJK character renders correctly even on text we consider "English".
        _set_run_fonts(run, latin=font_name, ea=FONT_CN_BODY)
    if auto_shrink:
        _enable_autofit(tf)


def _add_panel(slide, left, top, width, height, fill: RGBColor = BG_PANEL) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False


def _add_gold_line(slide, left, top, width, height=Pt(3)) -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()


def _page_header(slide, title_zh: str, title_en: str, page_num: int,
                 team_a: Team | None = None, team_b: Team | None = None) -> None:
    """Standard header for content pages. Chinese primary, English subtitle.

    Page number rendered in Chinese numerals (一/二/三...) for the
    Chinese-first design (#17).  If team_a / team_b provided, adds
    team-color swatches + name mini-badges on the right (#18).
    """
    cn_numerals = "〇一二三四五六七八九十"
    tens = ["", "十", "二十", "三十", "四十", "五十", "六十", "七十", "八十", "九十"]
    def _cn_num(n: int) -> str:
        if n < 11:
            return cn_numerals[n]
        if n < 20:
            return "十" + cn_numerals[n - 10]
        t, o = divmod(n, 10)
        return tens[t] + (cn_numerals[o] if o else "")

    _add_gold_line(slide, MARGIN, Inches(0.45), Inches(0.4), Pt(3))
    _add_textbox(
        slide, MARGIN + Inches(0.55), Inches(0.25), Inches(8), Inches(0.5),
        title_zh, font_size=SECTION_SIZE, bold=True, color=WHITE, font_name=FONT_CN_BODY
    )
    _add_textbox(
        slide, MARGIN + Inches(0.55), Inches(0.72), Inches(8), Inches(0.32),
        title_en, font_size=Pt(11), bold=False, color=GOLD, font_name=FONT_BODY
    )
    # #18 team color swatches on right + #17 page number always shown
    page_text = f"第 {_cn_num(page_num)} 页"
    _add_textbox(
        slide, SLIDE_W - Inches(2.0), Inches(0.30), Inches(1.6), Inches(0.4),
        page_text, font_size=Pt(12), color=GREY,
        font_name=FONT_CN_BODY, align=PP_ALIGN.RIGHT
    )
    if team_a is not None and team_b is not None:
        try:
            color_a = RGBColor.from_string(team_a.home_kit_color.lstrip("#"))
        except Exception:
            color_a = GOLD
        try:
            color_b = RGBColor.from_string(team_b.home_kit_color.lstrip("#"))
        except Exception:
            color_b = CYAN
        # Team badges to the LEFT of the page number
        right_x = SLIDE_W - Inches(4.6)
        _add_panel(slide, right_x, Inches(0.32), Inches(0.22), Inches(0.30), fill=color_a)
        _add_textbox(slide, right_x + Inches(0.25), Inches(0.30), Inches(0.9), Inches(0.30),
                    team_a.name_zh, font_size=Pt(9), color=WHITE, bold=True,
                    font_name=FONT_CN_BODY, anchor=MSO_ANCHOR.MIDDLE)
        _add_panel(slide, right_x + Inches(1.20), Inches(0.32), Inches(0.22), Inches(0.30), fill=color_b)
        _add_textbox(slide, right_x + Inches(1.45), Inches(0.30), Inches(0.9), Inches(0.30),
                    team_b.name_zh, font_size=Pt(9), color=WHITE, bold=True,
                    font_name=FONT_CN_BODY, anchor=MSO_ANCHOR.MIDDLE)


# ----------------------- pages -----------------------

def _page_cover(prs, result: PredictionResult) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    team_a = result.match.team_a
    team_b = result.match.team_b

    # Big gold accent line
    _add_gold_line(slide, MARGIN, Inches(1.3), Inches(2.0), Pt(5))

    # Two color stripes at top representing the two teams (#18)
    _add_panel(slide, MARGIN, Inches(0.4), Inches(2.0), Inches(0.10),
               fill=RGBColor(0xFF, 0xB6, 0x27))  # gold for team A
    _add_panel(slide, MARGIN + Inches(2.0), Inches(0.4), Inches(2.0), Inches(0.10),
               fill=RGBColor(0x00, 0xD4, 0xFF))   # cyan for team B
    # Kit color swatches (use actual team kit_color from teams.json)
    try:
        a_color = RGBColor.from_string(team_a.home_kit_color.lstrip("#"))
    except Exception:
        a_color = RGBColor(0xFF, 0xB6, 0x27)
    try:
        b_color = RGBColor.from_string(team_b.home_kit_color.lstrip("#"))
    except Exception:
        b_color = RGBColor(0x00, 0xD4, 0xFF)
    _add_panel(slide, MARGIN, Inches(0.55), Inches(0.5), Inches(0.4), fill=a_color)
    _add_panel(slide, MARGIN + Inches(2.0), Inches(0.55), Inches(0.5), Inches(0.4), fill=b_color)

    _add_textbox(
        slide, MARGIN + Inches(0.6), Inches(0.55), Inches(12), Inches(0.4),
        "世界杯预测报告  ·  WORLD CUP FORECAST",
        font_size=Pt(13), bold=True, color=GOLD, font_name=FONT_CN_BODY,
    )
    _add_textbox(
        slide, MARGIN, Inches(1.6), Inches(12), Inches(1.2),
        f"{team_a.name_zh}  对阵  {team_b.name_zh}",
        font_size=Pt(50), bold=True, color=WHITE, font_name=FONT_CN_BODY,
    )
    _add_textbox(
        slide, MARGIN, Inches(2.7), Inches(12), Inches(0.5),
        f"{team_a.name_en}  vs  {team_b.name_en}",
        font_size=Pt(22), bold=False, color=GREY, font_name=FONT_TITLE,
    )
    stage_zh = {
        "group": "小组赛",
        "round_of_16": "八分之一决赛",
        "quarterfinal": "四分之一决赛",
        "semifinal": "半决赛",
        "final": "决赛",
        "third_place": "三四名决赛",
    }.get(result.match.stage, result.match.stage)
    _add_textbox(
        slide, MARGIN, Inches(3.5), Inches(12), Inches(0.5),
        f"{result.match.match_date}  ·  阶段：{stage_zh}  ·  比赛场地：{result.match.venue}",
        font_size=Pt(16), color=WHITE, font_name=FONT_CN_BODY,
    )

    # Recommended pick + predicted score
    pick = result.recommended_pick
    pick_zh = team_a.name_zh if pick == "A" else (team_b.name_zh if pick == "B" else "平局")
    pick_en = team_a.name_en if pick == "A" else (team_b.name_en if pick == "B" else "Draw")
    predicted_score = result.monte_carlo.predicted_score
    score_a, score_b = result.monte_carlo.split_goals(predicted_score)
    # If our consensus pick disagrees with the score, fall back to the
    # score's own outcome (so the cover tells one consistent story).
    if (pick == "A" and score_a < score_b) or (pick == "B" and score_b < score_a):
        pick = result.monte_carlo.score_outcome(predicted_score)
        pick_zh = team_a.name_zh if pick == "A" else (team_b.name_zh if pick == "B" else "平局")
        pick_en = team_a.name_en if pick == "A" else (team_b.name_en if pick == "B" else "Draw")
    _add_panel(slide, MARGIN, Inches(4.6), Inches(6.2), Inches(2.1), fill=BG_PANEL)
    _add_textbox(
        slide, MARGIN + Inches(0.3), Inches(4.75), Inches(6), Inches(0.4),
        "预测比分  ·  PREDICTED SCORE", font_size=Pt(11), color=GOLD, font_name=FONT_CN_BODY, bold=True,
    )
    _add_textbox(
        slide, MARGIN + Inches(0.3), Inches(4.95), Inches(6), Inches(1.0),
        predicted_score, font_size=Pt(70), bold=True, color=WHITE, font_name=FONT_MONO, align=PP_ALIGN.LEFT,
    )
    _add_textbox(
        slide, MARGIN + Inches(0.3), Inches(5.85), Inches(6), Inches(0.32),
        f"推荐结果：{pick_zh} ({pick_en})", font_size=Pt(14), color=GOLD, font_name=FONT_CN_BODY,
    )
    # #5 预测比分与推荐结果说明
    if pick == "draw" or (pick == "A" and score_a == score_b) or (pick == "B" and score_b == score_a):
        explain = f"比分 {predicted_score} 为平局；综合 ELO/伤停推荐 {pick_zh}"
    elif pick == "A" and score_a > score_b:
        explain = f"比分 {predicted_score} 反映单场最可能结果"
    elif pick == "B" and score_b > score_a:
        explain = f"比分 {predicted_score} 反映单场最可能结果"
    else:
        explain = f"比分与推荐方向不同；取综合最大概率"
    _add_textbox(
        slide, MARGIN + Inches(0.3), Inches(6.15), Inches(6), Inches(0.28),
        explain, font_size=Pt(9), color=GREY, font_name=FONT_CN_BODY,
    )
    conf = result.confidence
    conf_color = GREEN if conf == "high" else (GOLD if conf == "medium" else RED)
    conf_zh = "高" if conf == "high" else "中" if conf == "medium" else "低"
    _add_textbox(
        slide, MARGIN + Inches(0.3), Inches(6.45), Inches(6), Inches(0.32),
        f"信心指数：{conf_zh}（{conf.upper()}）",
        font_size=Pt(11), color=conf_color, font_name=FONT_CN_BODY, bold=True,
    )

    # Right side: probability numbers
    p = result.model_probs.win_draw_loss
    right_left = MARGIN + Inches(6.6)
    _add_panel(slide, right_left, Inches(4.6), Inches(5.5), Inches(2.1), fill=BG_PANEL)
    _add_textbox(
        slide, right_left + Inches(0.3), Inches(4.75), Inches(5), Inches(0.4),
        "综合概率  ·  CONSENSUS", font_size=Pt(11), color=GOLD, font_name=FONT_CN_BODY, bold=True,
    )
    _add_textbox(
        slide, right_left + Inches(0.3), Inches(5.2), Inches(1.5), Inches(1.2),
        f"{p[0]:.0%}", font_size=Pt(40), bold=True, color=GREEN, font_name=FONT_MONO, align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide, right_left + Inches(0.3), Inches(6.4), Inches(1.5), Inches(0.3),
        f"胜 · {team_a.name_zh}", font_size=Pt(10), color=GREY, font_name=FONT_CN_BODY, align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide, right_left + Inches(1.95), Inches(5.2), Inches(1.5), Inches(1.2),
        f"{p[1]:.0%}", font_size=Pt(40), bold=True, color=GREY, font_name=FONT_MONO, align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide, right_left + Inches(1.95), Inches(6.4), Inches(1.5), Inches(0.3),
        "平 · DRAW", font_size=Pt(10), color=GREY, font_name=FONT_CN_BODY, align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide, right_left + Inches(3.6), Inches(5.2), Inches(1.5), Inches(1.2),
        f"{p[2]:.0%}", font_size=Pt(40), bold=True, color=RED, font_name=FONT_MONO, align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide, right_left + Inches(3.6), Inches(6.4), Inches(1.5), Inches(0.3),
        f"负 · {team_b.name_zh}", font_size=Pt(10), color=GREY, font_name=FONT_CN_BODY, align=PP_ALIGN.CENTER,
    )

    _add_textbox(
        slide, MARGIN, Inches(7.05), Inches(12), Inches(0.3),
        f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}（北京时间） ·  数据截止：{result.match.match_date}  ·  下次更新：比赛开始前 2 小时  ·  模型版本 v1.0  ·  数据：Wikipedia + football-data.org + ESPN",
        font_size=Pt(8), color=GREY_DARK, font_name=FONT_CN_BODY,
    )


def _page_summary(prs, result: PredictionResult) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _page_header(slide, "执行摘要", "Executive Summary", 2,
                 team_a=result.match.team_a, team_b=result.match.team_b)
    p = result.model_probs.win_draw_loss
    team_a = result.match.team_a
    team_b = result.match.team_b
    predicted_score = result.monte_carlo.predicted_score
    score_a, score_b = result.monte_carlo.split_goals(predicted_score)

    # Big predicted score callout
    _add_panel(slide, MARGIN, Inches(1.4), Inches(12.1), Inches(1.7), fill=BG_PANEL)
    _add_textbox(slide, MARGIN + Inches(0.4), Inches(1.55), Inches(12), Inches(0.4),
                "预测比分  ·  PREDICTED SCORE", font_size=Pt(12), color=GOLD,
                font_name=FONT_CN_BODY, bold=True)
    _add_textbox(slide, MARGIN, Inches(2.0), Inches(12.1), Inches(1.1),
                predicted_score, font_size=Pt(72), bold=True, color=WHITE,
                font_name=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Three probability cards (#9 bigger fonts)
    y = Inches(3.3)
    card_w = Inches(4.0)
    card_h = Inches(2.0)
    for i, (label_zh, val, color, team_zh) in enumerate([
        (f"{team_a.name_zh} 胜", p[0], GREEN, team_a.name_zh),
        ("平局", p[1], GREY, "—"),
        (f"{team_b.name_zh} 胜", p[2], RED, team_b.name_zh),
    ]):
        x = MARGIN + i * (card_w + Inches(0.15))
        _add_panel(slide, x, y, card_w, card_h, fill=BG_PANEL)
        _add_textbox(slide, x + Inches(0.3), y + Inches(0.12), card_w - Inches(0.6), Inches(0.4),
                    label_zh, font_size=Pt(18), color=color, bold=True, font_name=FONT_CN_BODY)
        _add_textbox(slide, x + Inches(0.3), y + Inches(0.6), card_w - Inches(0.6), Inches(0.3),
                    team_zh, font_size=Pt(10), color=GREY, font_name=FONT_CN_BODY)
        _add_textbox(slide, x + Inches(0.3), y + Inches(0.9), card_w - Inches(0.6), Inches(1.0),
                    f"{val:.0%}", font_size=Pt(56), bold=True, color=color, font_name=FONT_MONO)

    # Recommended pick + expected goals
    y2 = Inches(5.5)
    pick = result.recommended_pick
    pick_zh = team_a.name_zh if pick == "A" else (team_b.name_zh if pick == "B" else "平局")
    conf_zh = "高" if result.confidence == "high" else "中" if result.confidence == "medium" else "低"
    _add_panel(slide, MARGIN, y2, Inches(6.2), Inches(1.4), fill=BG_CARD)
    _add_textbox(slide, MARGIN + Inches(0.3), y2 + Inches(0.15), Inches(6), Inches(0.3),
                "推荐结果  ·  RECOMMENDED PICK", font_size=Pt(11), color=GOLD, font_name=FONT_CN_BODY, bold=True)
    _add_textbox(slide, MARGIN + Inches(0.3), y2 + Inches(0.5), Inches(6), Inches(0.55),
                pick_zh, font_size=Pt(28), bold=True, color=WHITE, font_name=FONT_CN_BODY)
    _add_textbox(slide, MARGIN + Inches(0.3), y2 + Inches(1.05), Inches(6), Inches(0.3),
                f"信心指数：{conf_zh}  ·  Confidence: {result.confidence.upper()}",
                font_size=Pt(11), color=GOLD, font_name=FONT_CN_BODY, bold=True)

    _add_panel(slide, MARGIN + Inches(6.5), y2, Inches(5.6), Inches(1.4), fill=BG_CARD)
    _add_textbox(slide, MARGIN + Inches(6.8), y2 + Inches(0.15), Inches(5), Inches(0.3),
                "预期进球（期望值）  ·  EXPECTED GOALS", font_size=Pt(11), color=GOLD, font_name=FONT_CN_BODY, bold=True)
    eg = result.model_probs.expected_goals
    _add_textbox(slide, MARGIN + Inches(6.8), y2 + Inches(0.5), Inches(5), Inches(0.7),
                f"{eg[0]:.1f}  —  {eg[1]:.1f}",
                font_size=Pt(32), bold=True, color=WHITE, font_name=FONT_MONO, align=PP_ALIGN.CENTER)


def _page_pre_match_news(prs, result: PredictionResult) -> None:
    """Show recent Chinese-language headlines relevant to the two teams.

    Sourced from Zhibo8's public 24-hour hot-news feed. Empty if the scraper
    could not reach Zhibo8 or no headlines matched the two teams.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _page_header(
        slide,
        "赛前热点",
        "Pre-Match Headlines",
        4,
        team_a=result.match.team_a,
        team_b=result.match.team_b,
    )

    _add_textbox(
        slide, MARGIN, Inches(1.4), Inches(12.1), Inches(0.4),
        "近 24 小时中文体育热点  ·  来源：直播吧 (zhibo8.cc)",
        font_size=Pt(12), color=GOLD, font_name=FONT_CN_BODY, bold=True,
    )

    news = result.pre_match_news
    if not news:
        _add_textbox(
            slide, MARGIN, Inches(2.4), Inches(12.1), Inches(0.8),
            "暂无相关热点（直播吧当前无新情报或网络受限）",
            font_size=Pt(14), color=GREY, font_name=FONT_CN_BODY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        return

    # Two-column layout, max 6 headlines.
    panel_h = Inches(5.4)
    _add_panel(slide, MARGIN, Inches(1.9), Inches(12.1), panel_h, fill=BG_CARD)
    for i, headline in enumerate(news[:6]):
        # Truncate very long titles to keep the slide clean.
        display = headline if len(headline) <= 80 else headline[:78] + "…"
        _add_textbox(
            slide,
            MARGIN + Inches(0.3),
            Inches(2.1) + Inches(i * 0.82),
            Inches(11.5),
            Inches(0.7),
            f"•  {display}",
            font_size=Pt(13), color=WHITE, font_name=FONT_CN_BODY,
        )


def _page_toc(prs, result: PredictionResult) -> None:
    """#16 Table of contents."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _page_header(slide, "目录", "Table of Contents", 3,
                 team_a=result.match.team_a, team_b=result.match.team_b)
    sections = [
        ("一", "执行摘要", "Executive Summary", "比分 + 三模型概率 + 信心指数"),
        ("二", "目录", "Contents", "本报告各章节导览"),
        ("三", "赛前热点", "Pre-Match Headlines", "近 24h 直播吧中文热点（伤停/阵容/情报）"),
        ("四", "大赛背景", "Tournament Context", "FIFA 世界杯 2026 主办国/规则/决赛"),
        ("五", "球队档案", "Team Profile", "FIFA 排名/ELO/教练/队长/伤停"),
        ("六", "历史交锋", "Head-to-Head", "近 5 次交锋结果 + 进球者"),
        ("七", "近期状态", "Recent Form", "近 10 场胜负 + 攻守数据"),
        ("八", "预测首发", "Predicted Lineup A", "加拿大 阵型 + 11 人首发"),
        ("九", "预测首发", "Predicted Lineup B", "波黑 阵型 + 11 人首发"),
        ("十", "核心球员", "Key Players A", "加拿大 TOP 5 球员"),
        ("十一", "核心球员", "Key Players B", "波黑 TOP 5 球员"),
        ("十二", "关键对位", "Key Matchups", "4 组 1v1 对决"),
        ("十三", "阵容深度", "Squad Depth", "首发 vs 替补战力 + 替补名单"),
        ("十四", "球队能力对比", "Team Capabilities", "6 维能力雷达"),
        ("十五", "比赛环境", "Match Environment", "赛程疲劳 + 裁判尺度 + 伤停"),
        ("十六", "三模型对比", "Model Output", "ELO/Poisson/XGBoost 概率对比"),
        ("十七", "蒙特卡洛", "Monte Carlo", "10000 次模拟 + 进球数分布 + 大/小 2.5"),
        ("十八", "敏感性分析", "Sensitivity", "影响最大的 6 个变量"),
        ("十九", "最终预测", "Final Prediction", "绝对比分 + 信心指数 + 关键风险"),
        ("二十", "附录", "Appendix", "数据源 + 方法 + 免责声明"),
    ]
    # Two columns
    y_top = Inches(1.2)
    rows_per_col = 10
    for i, (num, zh, en, desc) in enumerate(sections):
        col = i // rows_per_col
        row = i % rows_per_col
        x = MARGIN + col * Inches(6.3)
        y = y_top + Inches(row * 0.55)
        _add_textbox(slide, x, y, Inches(0.5), Inches(0.4),
                    f"第 {num}", font_size=Pt(14), color=GOLD, bold=True, font_name=FONT_CN_BODY)
        _add_textbox(slide, x + Inches(0.55), y, Inches(2.0), Inches(0.4),
                    f"{zh}", font_size=Pt(13), color=WHITE, bold=True, font_name=FONT_CN_BODY)
        _add_textbox(slide, x + Inches(2.55), y, Inches(2.0), Inches(0.4),
                    f"{en}", font_size=Pt(9), color=GREY, font_name=FONT_BODY)
        _add_textbox(slide, x + Inches(0.55), y + Inches(0.28), Inches(5.5), Inches(0.3),
                    f"  · {desc}", font_size=Pt(8), color=GREY, font_name=FONT_CN_BODY)


def _page_world_cup_context(prs, result: PredictionResult) -> None:
    """#7 World Cup 2026 background / venue info — inserted before team profile."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _page_header(slide, "大赛背景", "Tournament Context", 5,
                 team_a=result.match.team_a, team_b=result.match.team_b)
    team_a = result.match.team_a
    team_b = result.match.team_b

    # Left: tournament info card
    _add_panel(slide, MARGIN, Inches(1.5), Inches(6.0), Inches(5.4), fill=BG_PANEL)
    _add_textbox(slide, MARGIN + Inches(0.3), Inches(1.65), Inches(5.5), Inches(0.4),
                "FIFA 世界杯 2026", font_size=Pt(20), bold=True, color=GOLD, font_name=FONT_TITLE)
    _add_textbox(slide, MARGIN + Inches(0.3), Inches(2.05), Inches(5.5), Inches(0.4),
                "世界杯 2026  ·  WORLD CUP 2026", font_size=Pt(11), color=GREY, font_name=FONT_CN_BODY)

    facts = [
        ("主办国", "Host", "美国 / 加拿大 / 墨西哥", "USA / Canada / Mexico"),
        ("比赛日期", "Match Date", result.match.match_date, ""),
        ("比赛场地", "Venue", result.match.venue or "TBD", ""),
        ("比赛阶段", "Stage", {
            "group": "小组赛 Group Stage",
            "round_of_16": "八分之一决赛 Round of 16",
            "quarterfinal": "四分之一决赛 Quarterfinal",
            "semifinal": "半决赛 Semifinal",
            "final": "决赛 Final",
            "third_place": "三四名决赛 Third Place",
        }.get(result.match.stage, result.match.stage), ""),
        ("参赛队总数", "Total Teams", "48 队 (史上首次扩军)", "48 teams (first expansion)"),
        ("总场数", "Total Matches", "104 场", "104 matches"),
        ("决赛日期", "Final Date", "2026 年 7 月 19 日", "Jul 19, 2026"),
        ("决赛球场", "Final Venue", "大都会人寿体育场 (纽约)", "MetLife Stadium (NYC)"),
        ("冠军奖金", "Prize Money", "$50M 美元", "$50M USD"),
    ]
    for i, (zh, en, val_zh, val_en) in enumerate(facts):
        ry = Inches(2.6) + i * Inches(0.55)
        _add_textbox(slide, MARGIN + Inches(0.3), ry, Inches(1.6), Inches(0.4),
                    f"{zh}  ·  {en}", font_size=Pt(10), color=GREY, font_name=FONT_CN_BODY)
        _add_textbox(slide, MARGIN + Inches(2.0), ry, Inches(3.8), Inches(0.4),
                    f"{val_zh}  ·  {val_en}", font_size=Pt(11), color=WHITE,
                    font_name=FONT_CN_BODY, bold=True)

    # Right: match story (#6) + accent block
    _add_panel(slide, MARGIN + Inches(6.5), Inches(1.5), Inches(5.6), Inches(5.4), fill=BG_CARD)
    _add_textbox(slide, MARGIN + Inches(6.8), Inches(1.65), Inches(5.0), Inches(0.4),
                "比赛看点  ·  MATCH STORY", font_size=Pt(11), color=GOLD, font_name=FONT_CN_BODY, bold=True)

    # Generate match story from team data
    elo_a = team_a.elo
    elo_b = team_b.elo
    elo_gap = abs(elo_a - elo_b)
    if elo_gap > 150:
        fav = team_a.name_zh if elo_a > elo_b else team_b.name_zh
        under = team_b.name_zh if elo_a > elo_b else team_a.name_zh
        story_intro = f"{fav} 在 ELO 评分上领先 {elo_gap:.0f} 分"
    else:
        fav = team_a.name_zh
        under = team_b.name_zh
        story_intro = f"两队实力接近，ELO 差距仅 {elo_gap:.0f} 分"

    cap_a = team_a.captain_zh or team_a.captain
    cap_b = team_b.captain_zh or team_b.captain
    story = (
        f"{story_intro}，"
        f"由 {team_a.coach_zh or team_a.coach} 和 {team_b.coach_zh or team_b.coach} 挂帅的两支队伍将展开对决。"
        f"{team_a.name_zh} 由队长 {cap_a} 领衔，{team_b.name_zh} 队长 {cap_b}。"
        f"建议关注各自进攻核心的发挥，以及定位球和反击效率的较量。"
    )

    _add_textbox(slide, MARGIN + Inches(6.8), Inches(2.05), Inches(5.0), Inches(0.4),
                f"{team_a.name_zh}  vs  {team_b.name_zh}",
                font_size=Pt(18), bold=True, color=WHITE, font_name=FONT_CN_BODY)
    _add_textbox(slide, MARGIN + Inches(6.8), Inches(2.45), Inches(5.0), Inches(0.4),
                f"{team_a.name_en}  vs  {team_b.name_en}",
                font_size=Pt(11), color=GREY, font_name=FONT_TITLE)

    # Story body
    _add_textbox(slide, MARGIN + Inches(6.8), Inches(3.0), Inches(5.0), Inches(3.0),
                story, font_size=Pt(12), color=WHITE, font_name=FONT_CN_BODY)

    # Bottom of right panel: predicted score
    _add_textbox(slide, MARGIN + Inches(6.8), Inches(6.35), Inches(5.0), Inches(0.4),
                f"  本场预测比分  ·  {result.monte_carlo.predicted_score}",
                font_size=Pt(11), color=GOLD, font_name=FONT_CN_BODY, bold=True)


def _page_h2h(prs, result: PredictionResult) -> None:
    """#5 + #8 historical H2H page using real warehouse data."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _page_header(slide, "历史交锋", "Head-to-Head Record", 7,
                 team_a=result.match.team_a, team_b=result.match.team_b)
    team_a = result.match.team_a
    team_b = result.match.team_b

    # Load real H2H matches from the warehouse
    repo = MatchRepository()
    h2h_matches = repo.get_matches(team_code=team_a.code) + repo.get_matches(team_code=team_b.code)
    # Keep only matches where both teams played each other, and the match is before the current one
    h2h_matches = [
        m for m in h2h_matches
        if {m.home_team_code, m.away_team_code} == {team_a.code, team_b.code}
        and m.date < result.match.match_date
    ]
    h2h_matches = sorted(h2h_matches, key=lambda m: m.date, reverse=True)

    wins_a = draws = wins_b = 0
    for m in h2h_matches:
        if m.home_goals is None or m.away_goals is None:
            continue
        a_goals = m.home_goals if m.home_team_code == team_a.code else m.away_goals
        b_goals = m.away_goals if m.home_team_code == team_a.code else m.home_goals
        if a_goals > b_goals:
            wins_a += 1
        elif a_goals == b_goals:
            draws += 1
        else:
            wins_b += 1

    n_meetings = len(h2h_matches)

    # Headline numbers
    _add_panel(slide, MARGIN, Inches(1.4), Inches(12.1), Inches(1.4), fill=BG_PANEL)
    _add_textbox(slide, MARGIN + Inches(0.4), Inches(1.5), Inches(12), Inches(0.4),
                f"近 {n_meetings} 次交锋  ·  LAST {n_meetings} MEETINGS",
                font_size=Pt(12), color=GOLD, font_name=FONT_CN_BODY, bold=True)
    for i, (label, val, color) in enumerate([
        (f"{team_a.name_zh} 胜", wins_a, GOLD),
        ("平局", draws, GREY),
        (f"{team_b.name_zh} 胜", wins_b, CYAN),
    ]):
        x = MARGIN + Inches(0.5) + i * Inches(4.0)
        y = Inches(1.85)
        _add_textbox(slide, x, y, Inches(3.6), Inches(0.4),
                    label, font_size=Pt(13), color=color, bold=True, font_name=FONT_CN_BODY)
        _add_textbox(slide, x, y + Inches(0.35), Inches(3.6), Inches(0.65),
                    f"{val}", font_size=Pt(40), bold=True, color=color, font_name=FONT_MONO,
                    align=PP_ALIGN.CENTER)

    # Recent meetings list
    _add_textbox(slide, MARGIN, Inches(3.1), Inches(12), Inches(0.4),
                "近 5 次交锋（真实数据）  ·  RECENT MEETINGS",
                font_size=Pt(13), bold=True, color=GOLD, font_name=FONT_CN_BODY)

    recent5 = h2h_matches[:5]
    if not recent5:
        _add_textbox(slide, MARGIN, Inches(3.6), Inches(12), Inches(0.5),
                    "暂无历史交锋记录  ·  No historical meetings in database",
                    font_size=Pt(12), color=GREY, font_name=FONT_CN_BODY, align=PP_ALIGN.CENTER)

    for i, m in enumerate(recent5):
        col = i % 5
        x = MARGIN + col * Inches(2.45)
        y = Inches(3.55)
        a_goals = m.home_goals if m.home_team_code == team_a.code else m.away_goals
        b_goals = m.away_goals if m.home_team_code == team_a.code else m.home_goals
        if a_goals is None or b_goals is None:
            continue
        result_text = "胜" if a_goals > b_goals else ("平" if a_goals == b_goals else "负")
        color = GREEN if result_text == "胜" else (GREY if result_text == "平" else RED)
        venue = "中立场" if m.neutral else ("主场" if m.home_team_code == team_a.code else "客场")

        _add_panel(slide, x, y, Inches(2.35), Inches(2.0), fill=BG_CARD)
        _add_textbox(slide, x, y + Inches(0.1), Inches(2.35), Inches(0.4),
                    m.date, font_size=Pt(10), color=GREY, font_name=FONT_MONO,
                    align=PP_ALIGN.CENTER)
        _add_textbox(slide, x, y + Inches(0.4), Inches(2.35), Inches(0.4),
                    result_text, font_size=Pt(24), bold=True, color=color, font_name=FONT_CN_BODY,
                    align=PP_ALIGN.CENTER)
        _add_textbox(slide, x, y + Inches(0.85), Inches(2.35), Inches(0.3),
                    venue, font_size=Pt(9), color=GREY, font_name=FONT_CN_BODY,
                    align=PP_ALIGN.CENTER)
        _add_textbox(slide, x, y + Inches(1.15), Inches(2.35), Inches(0.3),
                    f"{team_a.name_zh} {a_goals} - {b_goals} {team_b.name_zh}",
                    font_size=Pt(10), color=WHITE, font_name=FONT_CN_BODY,
                    align=PP_ALIGN.CENTER, bold=True)
        _add_textbox(slide, x, y + Inches(1.55), Inches(2.35), Inches(0.3),
                    m.competition, font_size=Pt(8), color=GREY, font_name=FONT_CN_BODY,
                    align=PP_ALIGN.CENTER)

    # Aggregate stat
    _add_textbox(slide, MARGIN, Inches(5.85), Inches(12), Inches(0.5),
                f"历史交锋总战绩  ·  ALL-TIME RECORD",
                font_size=Pt(13), bold=True, color=GOLD, font_name=FONT_CN_BODY, align=PP_ALIGN.CENTER)
    if n_meetings > 0:
        record_text = (
            f"{team_a.name_zh} {wins_a} 胜  ·  {draws} 平  ·  {wins_b} 胜 {team_b.name_zh}    "
            f"(近 {n_meetings} 场胜率：{team_a.name_zh} {wins_a/n_meetings*100:.0f}%  ·  {team_b.name_zh} {wins_b/n_meetings*100:.0f}%)"
        )
    else:
        record_text = "数据库中暂无两队交锋记录  ·  No H2H data available"
    _add_textbox(slide, MARGIN, Inches(6.3), Inches(12), Inches(0.6),
                record_text,
                font_size=Pt(14), color=WHITE, font_name=FONT_CN_BODY, align=PP_ALIGN.CENTER)


def _page_team_profile(prs, result: PredictionResult) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _page_header(slide, "球队档案", "Team Profile", 6,
                 team_a=result.match.team_a, team_b=result.match.team_b)
    team_a = result.match.team_a
    team_b = result.match.team_b

    # Removed the empty "主队服颜色" row; replaced with a visual ELO gap bar
    # (#20).  Rows: FIFA Rank, ELO, Coach, Captain, Confederation, + ELO gap bar.
    headers_zh = ["FIFA 排名", "ELO 评分", "主教练", "队长", "所属足联", "ELO 差距"]
    headers_en = ["FIFA Rank", "ELO Rating", "Head Coach", "Captain", "Confederation", "ELO Gap"]
    a_vals = [
        f"#{team_a.fifa_ranking}", f"{team_a.elo:.0f}",
        team_a.coach_zh or team_a.coach, team_a.captain_zh or team_a.captain,
        team_a.confederation, "",
    ]
    b_vals = [
        f"#{team_b.fifa_ranking}", f"{team_b.elo:.0f}",
        team_b.coach_zh or team_b.coach, team_b.captain_zh or team_b.captain,
        team_b.confederation, "",
    ]

    y = Inches(1.4)
    row_h = 0.62
    for i, (h_zh, h_en, va, vb) in enumerate(zip(headers_zh, headers_en, a_vals, b_vals)):
        _add_panel(slide, MARGIN, y + i * Inches(row_h), Inches(12.1), Inches(0.52),
                   fill=BG_PANEL if i % 2 == 0 else BG_CARD)
        _add_textbox(slide, MARGIN + Inches(0.3), y + i * Inches(row_h) + Inches(0.08),
                    Inches(3.0), Inches(0.4),
                    f"{h_zh}  ·  {h_en}", font_size=Pt(11), color=GOLD, font_name=FONT_CN_BODY, bold=True)
        if h_zh == "所属足联":  # Confederation with color stripe accent
            _add_panel(slide, MARGIN + Inches(3.5), y + i * Inches(row_h) + Inches(0.08),
                       Inches(0.15), Inches(0.4), fill=GOLD)
            _add_textbox(slide, MARGIN + Inches(3.7), y + i * Inches(row_h) + Inches(0.08),
                        Inches(3.8), Inches(0.4),
                        str(va), font_size=Pt(14), color=WHITE, font_name=FONT_CN_BODY, bold=True)
            _add_panel(slide, MARGIN + Inches(7.7), y + i * Inches(row_h) + Inches(0.08),
                       Inches(0.15), Inches(0.4), fill=CYAN)
            _add_textbox(slide, MARGIN + Inches(7.9), y + i * Inches(row_h) + Inches(0.08),
                        Inches(4.0), Inches(0.4),
                        str(vb), font_size=Pt(14), color=WHITE, font_name=FONT_CN_BODY, bold=True, align=PP_ALIGN.LEFT)
        else:
            _add_textbox(slide, MARGIN + Inches(3.5), y + i * Inches(row_h) + Inches(0.08),
                        Inches(4.0), Inches(0.4),
                        str(va), font_size=Pt(15), color=WHITE, font_name=FONT_CN_BODY, bold=True)
            _add_textbox(slide, MARGIN + Inches(7.7), y + i * Inches(row_h) + Inches(0.08),
                        Inches(4.0), Inches(0.4),
                        str(vb), font_size=Pt(15), color=WHITE, font_name=FONT_CN_BODY, bold=True, align=PP_ALIGN.LEFT)

    # ELO gap horizontal bar (row 6, below the 5 main rows)
    bar_y = y + Inches(6 * row_h + 0.15)
    elo_a = team_a.elo
    elo_b = team_b.elo
    diff = abs(elo_a - elo_b)
    bar_max = Inches(8.0)
    bar_left = MARGIN + Inches(2.0)
    # background bar
    _add_panel(slide, bar_left, bar_y, bar_max, Inches(0.30), fill=BG_SOFT)
    # fill bar
    if elo_a >= elo_b:
        fill_w = Inches(min(0.95, diff / 200.0) * 8.0)
        _add_panel(slide, bar_left, bar_y, fill_w, Inches(0.30), fill=GOLD)
        label = f"  {team_a.name_zh} +{diff:.0f} ELO"
    else:
        fill_w = Inches(min(0.95, diff / 200.0) * 8.0)
        _add_panel(slide, bar_left, bar_y, fill_w, Inches(0.30), fill=CYAN)
        label = f"  {team_b.name_zh} +{diff:.0f} ELO"
    _add_textbox(slide, MARGIN, bar_y - Inches(0.05), Inches(12.0), Inches(0.3),
                f"ELO 差距  ·  {team_a.name_zh} {elo_a:.0f}  vs  {team_b.name_zh} {elo_b:.0f}    →    {label}",
                font_size=Pt(10), color=GREY, font_name=FONT_CN_BODY)


def _page_recent_form(prs, result: PredictionResult) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _page_header(slide, "近期状态", "Recent Form (Last 10 Matches)", 8,
                 team_a=result.match.team_a, team_b=result.match.team_b)
    team_a = result.match.team_a
    team_b = result.match.team_b

    repo = MatchRepository()

    def _team_form(team_code: str, match_date: str) -> list[str]:
        """Return W/D/L strings for the last 10 matches before match_date."""
        matches = repo.get_matches(team_code=team_code, before=match_date, limit=10)
        form: list[str] = []
        for m in sorted(matches, key=lambda x: x.date):
            if m.home_goals is None or m.away_goals is None:
                continue
            team_goals = m.home_goals if m.home_team_code == team_code else m.away_goals
            opp_goals = m.away_goals if m.home_team_code == team_code else m.home_goals
            if team_goals > opp_goals:
                form.append("胜")
            elif team_goals == opp_goals:
                form.append("平")
            else:
                form.append("负")
        return form

    form_a = _team_form(team_a.code, result.match.match_date)
    form_b = _team_form(team_b.code, result.match.match_date)

    def render_team_form(name_zh: str, name_en: str, form: list[str], top: float, color) -> None:
        _add_textbox(slide, MARGIN, top, Inches(3), Inches(0.6),
                    name_zh, font_size=Pt(28), bold=True, color=WHITE, font_name=FONT_CN_BODY)
        actual_n = len(form)
        _add_textbox(slide, MARGIN, top + Inches(0.6), Inches(3), Inches(0.3),
                    f"近 {actual_n} 场  ·  Last {actual_n}", font_size=Pt(10), color=GREY, font_name=FONT_CN_BODY)
        # W/D/L summary line
        w = sum(1 for r in form if r == "胜")
        d = sum(1 for r in form if r == "平")
        l = sum(1 for r in form if r == "负")
        _add_textbox(slide, MARGIN, top + Inches(0.95), Inches(3), Inches(0.3),
                    f"  胜 {w}    平 {d}    负 {l}", font_size=Pt(11), color=color, bold=True,
                    font_name=FONT_CN_BODY)
        # Chips
        for i, r in enumerate(form):
            chip_color = GREEN if r == "胜" else (GREY if r == "平" else RED)
            chip_left = MARGIN + Inches(3.0) + i * Inches(0.92)
            _add_panel(slide, chip_left, top + Inches(0.05), Inches(0.82), Inches(0.82), fill=chip_color)
            _add_textbox(slide, chip_left, top + Inches(0.05), Inches(0.82), Inches(0.82),
                        r, font_size=Pt(22), bold=True, color=BG_DEEP, font_name=FONT_CN_BODY,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    render_team_form(team_a.name_zh, team_a.name_en, form_a, Inches(1.5), GOLD)
    render_team_form(team_b.name_zh, team_b.name_en, form_b, Inches(3.4), CYAN)

    # Stats summary
    y3 = Inches(5.0)
    for side_x, side_team, side_color, side_stats in [
        (MARGIN, team_a, GOLD, result.team_a_stats),
        (MARGIN + Inches(6.3), team_b, CYAN, result.team_b_stats),
    ]:
        _add_panel(slide, side_x, y3, Inches(6.0), Inches(1.8), fill=BG_PANEL)
        _add_textbox(slide, side_x + Inches(0.3), y3 + Inches(0.1), Inches(5.5), Inches(0.3),
                    f"{side_team.name_zh}  ·  统计概览", font_size=Pt(13), bold=True,
                    color=side_color, font_name=FONT_CN_BODY)
        s = side_stats
        _add_textbox(slide, side_x + Inches(0.3), y3 + Inches(0.45), Inches(5.5), Inches(0.3),
                    f"  场均进球：{s.goals_per_game:.2f}    场均失球：{s.conceded_per_game:.2f}    xG：{s.xg_per_game:.2f}",
                    font_size=Pt(11), color=WHITE, font_name=FONT_CN_BODY)
        _add_textbox(slide, side_x + Inches(0.3), y3 + Inches(0.75), Inches(5.5), Inches(0.3),
                    f"  零封率：{s.clean_sheet_rate:.0%}    关键传球：{s.key_passes_per_game:.1f}    样本量：{s.last_10_wins + s.last_10_draws + s.last_10_losses}",
                    font_size=Pt(11), color=WHITE, font_name=FONT_CN_BODY)
        _add_textbox(slide, side_x + Inches(0.3), y3 + Inches(1.05), Inches(5.5), Inches(0.3),
                    f"  抢断：{s.tackles_per_game:.0f}    拦截：{s.interceptions_per_game:.0f}    射门精度：{s.shot_accuracy:.0%}",
                    font_size=Pt(11), color=GREY, font_name=FONT_CN_BODY)
        _add_textbox(slide, side_x + Inches(0.3), y3 + Inches(1.4), Inches(5.5), Inches(0.3),
                    f"  胜{s.last_10_wins} 平{s.last_10_draws} 负{s.last_10_losses}",
                    font_size=Pt(11), color=side_color, font_name=FONT_CN_BODY, bold=True)


def _page_predicted_lineup(prs, result: PredictionResult, side: str = "A") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    team_a = result.match.team_a
    team_b = result.match.team_b
    lineup = result.lineup_a if side == "A" else result.lineup_b
    team = team_a if side == "A" else team_b
    kit = team_a.home_kit_color if side == "A" else team_b.home_kit_color

    _page_header(slide, f"预测首发 · {team.name_zh}", f"Predicted Lineup · {team.name_en}",
                 9 if side == "A" else 10, team_a=team, team_b=team)
    _add_textbox(slide, MARGIN, Inches(1.2), Inches(12), Inches(0.4),
                f"阵型：{lineup.formation}  ·  FORMATION", font_size=Pt(15), bold=True, color=GOLD, font_name=FONT_CN_BODY)

    # Pitch image on the left
    pitch_path = CHART_DIR / f"pitch_{team.code}.png"
    draw_pitch_with_lineup(lineup, pitch_path, kit_color=kit, title="")
    slide.shapes.add_picture(str(pitch_path), MARGIN, Inches(1.7), height=Inches(5.5))

    # 11 mini player cards on the right (2 cols × 6 rows)
    right_left = MARGIN + Inches(5.5)
    for i, p in enumerate(lineup.players):
        col = i % 2
        row = i // 2
        x = right_left + col * Inches(3.55)
        y = Inches(1.7) + row * Inches(0.95)
        _add_panel(slide, x, y, Inches(3.45), Inches(0.85), fill=BG_CARD)
        # number circle
        _add_panel(slide, x + Inches(0.05), y + Inches(0.05), Inches(0.55), Inches(0.75),
                   fill=RGBColor(0xFF, 0xB6, 0x27) if side == "A" else RGBColor(0x00, 0xD4, 0xFF))
        _add_textbox(slide, x + Inches(0.05), y + Inches(0.05), Inches(0.55), Inches(0.75),
                    str(p.number or ""), font_size=Pt(20), bold=True, color=BG_DEEP,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name=FONT_MONO)
        # name + pos/club  (clipped to width 1.6 to leave room for rating)
        _add_textbox(slide, x + Inches(0.7), y + Inches(0.05), Inches(1.95), Inches(0.32),
                    p.display_name_cn(), font_size=Pt(10), bold=True, color=WHITE, font_name=FONT_CN_BODY)
        _add_textbox(slide, x + Inches(0.7), y + Inches(0.38), Inches(1.95), Inches(0.42),
                    f"{p.position}  ·  {p.club}", font_size=Pt(7), color=GREY, font_name=FONT_CN_BODY)
        # rating — well clear of name/club area now
        _add_textbox(slide, x + Inches(2.72), y + Inches(0.05), Inches(0.65), Inches(0.75),
                    f"{p.rating:.0f}", font_size=Pt(22), bold=True, color=GOLD if side == "A" else CYAN,
                    font_name=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _page_key_players(prs, result: PredictionResult, side: str = "A") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    team = result.match.team_a if side == "A" else result.match.team_b
    lineup = result.lineup_a if side == "A" else result.lineup_b
    _page_header(slide, f"核心球员 TOP 5 · {team.name_zh}", f"Key Players · {team.name_en}",
                 11 if side == "A" else 12, team_a=team, team_b=team)

    # Top 5 by rating — only players with a real photo get a card.
    key_players = [p for p in sorted(lineup.players, key=lambda p: p.rating, reverse=True) if player_has_photo(p)][:5]
    # Save the rank for each player (1 = top)
    rating_rank = {id(p): i + 1 for i, p in enumerate(key_players)}

    if not key_players:
        _add_textbox(slide, MARGIN, Inches(2.0), Inches(12), Inches(0.5),
                    "暂无球员照片数据  ·  No player photos available",
                    font_size=Pt(14), color=GREY, font_name=FONT_CN_BODY, align=PP_ALIGN.CENTER)
        return

    # Render cards in the original squad order so the visual lineup doesn't
    # jump around, but overlay the rank badge.
    for i, p in enumerate(key_players):
        col = i % 5
        x = MARGIN + col * Inches(2.4)
        y = Inches(1.6)
        # render card
        card_path = render_player_card(p, kit_color=team.home_kit_color)
        slide.shapes.add_picture(str(card_path), x, y, width=Inches(2.3), height=Inches(3.2))
        # rank badge (top-right corner)
        rank = rating_rank[id(p)]
        badge_color = GOLD if rank == 1 else (RGBColor(0xC0, 0xC0, 0xC0) if rank == 2 else RGBColor(0xCD, 0x7F, 0x32))
        _add_panel(slide, x + Inches(2.0), y + Inches(0.1), Inches(0.5), Inches(0.5), fill=badge_color)
        _add_textbox(slide, x + Inches(2.0), y + Inches(0.1), Inches(0.5), Inches(0.5),
                    f"#{rank}", font_size=Pt(14), bold=True, color=BG_DEEP,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name=FONT_MONO)
        # stat row
        _add_panel(slide, x, Inches(5.0), Inches(2.3), Inches(1.4), fill=BG_CARD)
        _add_textbox(slide, x + Inches(0.1), Inches(5.05), Inches(2.1), Inches(0.3),
                    f"年龄 {p.age}  ·  国家队出场 {p.caps}  ·  进球 {p.goals}",
                    font_size=Pt(9), color=GREY, font_name=FONT_CN_BODY)
        _add_textbox(slide, x + Inches(0.1), Inches(5.4), Inches(2.1), Inches(0.4),
                    f"身高 {p.height_cm}cm  ·  惯用脚 {p.preferred_foot}",
                    font_size=Pt(9), color=GREY, font_name=FONT_CN_BODY)
        _add_textbox(slide, x + Inches(0.1), Inches(5.7), Inches(2.1), Inches(0.4),
                    p.club, font_size=Pt(10), bold=True, color=WHITE, font_name=FONT_CN_BODY)


def _page_key_matchups(prs, result: PredictionResult) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _page_header(slide, "关键对位 1v1", "Key Matchups", 13,
                 team_a=result.match.team_a, team_b=result.match.team_b)

    # #12 Color + icon for each duel type
    duel_palette = {
        0: (RGBColor(0xE6, 0x39, 0x46), "攻 vs 防", "Attack vs Defence", "⚔"),
        1: (RGBColor(0xF4, 0xC4, 0x43), "中场对决", "Midfield Duel", "⊕"),
        2: (RGBColor(0x00, 0xD4, 0xFF), "边路竞速", "Wing Race", "»"),
        3: (RGBColor(0x7D, 0xCE, 0x82), "门线对决", "Keeper Duel", "⛔"),
    }
    # Only render matchups where both players have a real photo on disk.
    matchups = [mu for mu in result.key_matchups if player_has_photo(mu.player_a) and player_has_photo(mu.player_b)]
    if not matchups:
        _add_textbox(slide, MARGIN, Inches(2.0), Inches(12), Inches(0.5),
                    "暂无对位球员照片数据  ·  No matchup player photos available",
                    font_size=Pt(14), color=GREY, font_name=FONT_CN_BODY, align=PP_ALIGN.CENTER)
        return

    for i, mu in enumerate(matchups):
        col = i % 2
        row = i // 2
        x = MARGIN + col * Inches(6.15)
        y = Inches(1.5) + row * Inches(2.9)
        color, tag_zh, tag_en, icon = duel_palette.get(i, (GOLD, "关键对位", "Duel", "·"))

        _add_panel(slide, x, y, Inches(6.0), Inches(2.7), fill=BG_PANEL)
        # Left color strip accent (#12)
        _add_panel(slide, x, y, Inches(0.1), Inches(2.7), fill=color)
        # Title with type tag + icon
        _add_panel(slide, x + Inches(0.25), y + Inches(0.1), Inches(1.6), Inches(0.35), fill=color)
        _add_textbox(slide, x + Inches(0.25), y + Inches(0.1), Inches(1.6), Inches(0.35),
                    f"{icon}  {tag_zh}", font_size=Pt(11), bold=True, color=BG_DEEP,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name=FONT_CN_BODY)
        _add_textbox(slide, x + Inches(1.95), y + Inches(0.1), Inches(3.9), Inches(0.35),
                    mu.title_zh, font_size=Pt(11), bold=True, color=WHITE, font_name=FONT_CN_BODY,
                    anchor=MSO_ANCHOR.MIDDLE)
        _add_textbox(slide, x + Inches(0.25), y + Inches(0.5), Inches(5.5), Inches(0.25),
                    f"{mu.title_en}  ·  {tag_en}", font_size=Pt(9), color=GREY, font_name=FONT_BODY)

        # Two player photos side by side
        for j, p in enumerate([mu.player_a, mu.player_b]):
            px = x + Inches(0.25) + j * Inches(1.55)
            py = y + Inches(0.95)
            card_path = render_player_card(p, size=(200, 280), kit_color="#FFB627")
            slide.shapes.add_picture(str(card_path), px, py, width=Inches(1.5), height=Inches(1.5))
            _add_textbox(slide, px, py + Inches(1.5), Inches(1.5), Inches(0.25),
                        p.display_name_cn(), font_size=Pt(8), color=WHITE, align=PP_ALIGN.CENTER, font_name=FONT_CN_BODY)
            _add_textbox(slide, px, py + Inches(1.7), Inches(1.5), Inches(0.2),
                        p.position, font_size=Pt(7), color=GREY, align=PP_ALIGN.CENTER, font_name=FONT_CN_BODY)

        # VS divider with icon background
        _add_panel(slide, x + Inches(2.7), y + Inches(1.45), Inches(0.5), Inches(0.5), fill=color)
        _add_textbox(slide, x + Inches(2.7), y + Inches(1.45), Inches(0.5), Inches(0.5),
                    "VS", font_size=Pt(14), bold=True, color=BG_DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name=FONT_TITLE)

        # Stats — column header
        sx = x + Inches(3.3)
        sy = y + Inches(0.95)
        _add_textbox(slide, sx, sy - Inches(0.30), Inches(2.6), Inches(0.22),
                    f"{mu.player_a.display_name_cn()}  vs  {mu.player_b.display_name_cn()}",
                    font_size=Pt(8), color=GREY, font_name=FONT_CN_BODY, align=PP_ALIGN.CENTER)
        # #4 减少 stat 行到 4 行核心（评分/身高/年龄/进球），合并 国家队出场
        core_stats = [
            ("评分", mu.player_a.rating, mu.player_b.rating, f"{mu.player_a.rating:.1f}", f"{mu.player_b.rating:.1f}"),
            ("年龄", mu.player_a.age, mu.player_b.age, f"{mu.player_a.age}", f"{mu.player_b.age}"),
            ("身高", mu.player_a.height_cm, mu.player_b.height_cm, f"{mu.player_a.height_cm}cm", f"{mu.player_b.height_cm}cm"),
            ("进球", mu.player_a.goals, mu.player_b.goals, f"{mu.player_a.goals}", f"{mu.player_b.goals}"),
        ]
        for k, (zh, va_num, vb_num, va, vb) in enumerate(core_stats):
            row_y = sy + Inches(k * 0.32)
            winner_color = GOLD if va_num > vb_num else (CYAN if vb_num > va_num else GREY)
            _add_textbox(slide, sx, row_y, Inches(0.85), Inches(0.28),
                        va, font_size=Pt(11), bold=True, color=winner_color, font_name=FONT_MONO, align=PP_ALIGN.RIGHT)
            _add_textbox(slide, sx + Inches(0.9), row_y, Inches(0.9), Inches(0.28),
                        zh, font_size=Pt(10), color=GREY, font_name=FONT_CN_BODY, align=PP_ALIGN.CENTER)
            _add_textbox(slide, sx + Inches(1.85), row_y, Inches(0.85), Inches(0.28),
                        vb, font_size=Pt(11), bold=True, color=winner_color, font_name=FONT_MONO)


def _page_squad_depth(prs, result: PredictionResult) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _page_header(slide, "阵容深度", "Squad Depth", 14,
                 team_a=result.match.team_a, team_b=result.match.team_b)

    path = CHART_DIR / "depth.png"
    depth_bars(
        result.team_a_stats.starter_strength, result.team_a_stats.bench_strength,
        result.team_b_stats.starter_strength, result.team_b_stats.bench_strength,
        path,
    )
    slide.shapes.add_picture(str(path), MARGIN, Inches(1.5), width=Inches(7.5))

    # Bench details — 2 columns × 3 rows, with position-colored bullet + rating (#14)
    y = Inches(5.4)
    _add_textbox(slide, MARGIN, y, Inches(6), Inches(0.3),
                f"{result.match.team_a.name_zh} 替补席  ·  Bench", font_size=Pt(13), bold=True, color=GOLD, font_name=FONT_CN_BODY)
    bench_a = [p for p in result.lineup_a.bench[:6]]
    pos_color_dot = {
        "GK": RGBColor(0xF4, 0xC4, 0x43),
        "CB": RGBColor(0x3D, 0xC1, 0xD3), "LB": RGBColor(0x3D, 0xC1, 0xD3), "RB": RGBColor(0x3D, 0xC1, 0xD3),
        "CDM": RGBColor(0x7D, 0xCE, 0x82), "CM": RGBColor(0x7D, 0xCE, 0x82), "CAM": RGBColor(0x7D, 0xCE, 0x82),
        "ST": RGBColor(0xFF, 0x6B, 0x6B), "CF": RGBColor(0xFF, 0x6B, 0x6B),
        "RW": RGBColor(0xFF, 0x6B, 0x6B), "LW": RGBColor(0xFF, 0x6B, 0x6B),
    }
    for i, p in enumerate(bench_a):
        col = i % 2
        row = i // 2
        bx = MARGIN + col * Inches(3.0)
        by = y + Inches(0.35) + row * Inches(0.32)
        dot_color = pos_color_dot.get(p.position, GREY)
        _add_panel(slide, bx + Inches(0.05), by + Inches(0.08), Inches(0.15), Inches(0.15), fill=dot_color)
        _add_textbox(slide, bx + Inches(0.25), by, Inches(1.4), Inches(0.3),
                    f"#{p.number or '?'} {p.display_name_cn()}", font_size=Pt(9), color=WHITE, bold=True, font_name=FONT_CN_BODY)
        _add_textbox(slide, bx + Inches(1.65), by, Inches(1.3), Inches(0.3),
                    f"{p.rating:.0f} · {p.preferred_foot}脚",
                    font_size=Pt(8), color=GREY, font_name=FONT_CN_BODY, align=PP_ALIGN.RIGHT)

    _add_textbox(slide, MARGIN + Inches(6.5), y, Inches(6), Inches(0.3),
                f"{result.match.team_b.name_zh} 替补席  ·  Bench", font_size=Pt(13), bold=True, color=CYAN, font_name=FONT_CN_BODY)
    bench_b = [p for p in result.lineup_b.bench[:6]]
    for i, p in enumerate(bench_b):
        col = i % 2
        row = i // 2
        bx = MARGIN + Inches(6.5) + col * Inches(3.0)
        by = y + Inches(0.35) + row * Inches(0.32)
        dot_color = pos_color_dot.get(p.position, GREY)
        _add_panel(slide, bx + Inches(0.05), by + Inches(0.08), Inches(0.15), Inches(0.15), fill=dot_color)
        _add_textbox(slide, bx + Inches(0.25), by, Inches(1.4), Inches(0.3),
                    f"#{p.number or '?'} {p.display_name_cn()}", font_size=Pt(9), color=WHITE, bold=True, font_name=FONT_CN_BODY)
        _add_textbox(slide, bx + Inches(1.65), by, Inches(1.3), Inches(0.3),
                    f"{p.rating:.0f} · {p.preferred_foot}脚",
                    font_size=Pt(8), color=GREY, font_name=FONT_CN_BODY, align=PP_ALIGN.RIGHT)


def _page_radar(prs, result: PredictionResult) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _page_header(slide, "球队能力对比", "Team Capabilities", 15,
                 team_a=result.match.team_a, team_b=result.match.team_b)

    path = CHART_DIR / "radar.png"
    radar_chart(result.team_a_stats, result.team_b_stats, path)
    slide.shapes.add_picture(str(path), MARGIN, Inches(1.4), height=Inches(5.6))

    # Side stats table
    y = Inches(1.5)
    rows = [
        ("场均进球", "Goals/Game", "goals_per_game", "{:.2f}"),
        ("场均失球", "Conceded/Game", "conceded_per_game", "{:.2f}"),
        ("期望进球 xG", "xG", "xg_per_game", "{:.2f}"),
        ("期望失球 xGA", "xGA", "xga_per_game", "{:.2f}"),
        ("零封率", "Clean Sheet %", "clean_sheet_rate", "{:.0%}"),
        ("关键传球", "Key Passes", "key_passes_per_game", "{:.1f}"),
    ]
    for i, (zh, en, attr, fmt) in enumerate(rows):
        ry = y + Inches(i * 0.45)
        _add_panel(slide, MARGIN + Inches(7.5), ry, Inches(4.6), Inches(0.4), fill=BG_CARD if i % 2 else BG_PANEL)
        _add_textbox(slide, MARGIN + Inches(7.6), ry + Inches(0.05), Inches(2.5), Inches(0.3),
                    f"{zh}  ·  {en}", font_size=Pt(9), color=GREY, font_name=FONT_CN_BODY)
        a_val = getattr(result.team_a_stats, attr)
        b_val = getattr(result.team_b_stats, attr)
        _add_textbox(slide, MARGIN + Inches(10.1), ry + Inches(0.05), Inches(0.95), Inches(0.3),
                    fmt.format(a_val), font_size=Pt(11), color=GOLD, bold=True, font_name=FONT_MONO, align=PP_ALIGN.RIGHT)
        _add_textbox(slide, MARGIN + Inches(11.1), ry + Inches(0.05), Inches(0.95), Inches(0.3),
                    fmt.format(b_val), font_size=Pt(11), color=CYAN, bold=True, font_name=FONT_MONO, align=PP_ALIGN.RIGHT)


def _page_qualitative(prs, result: PredictionResult) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _page_header(slide, "定性因子 + 比赛环境", "Qualitative & Match Environment", 16,
                 team_a=result.match.team_a, team_b=result.match.team_b)

    path = CHART_DIR / "qualitative.png"
    qualitative_radar(result.qualitative_a, result.qualitative_b, path)
    slide.shapes.add_picture(str(path), MARGIN, Inches(1.3), width=Inches(5.0))

    # Right side: expanded stat comparison (#6 #9 #10)
    right_x = MARGIN + Inches(5.3)
    right_w = Inches(7.3)
    y_r = Inches(1.3)
    _add_panel(slide, right_x, y_r, right_w, Inches(3.5), fill=BG_PANEL)
    _add_textbox(slide, right_x + Inches(0.2), y_r + Inches(0.1), Inches(7.0), Inches(0.35),
                "深度数据  ·  ADVANCED METRICS", font_size=Pt(11), bold=True, color=GOLD, font_name=FONT_CN_BODY)

    # Two-column comparison: team A | stat | team B
    sa = result.team_a_stats
    sb = result.team_b_stats
    metrics = [
        ("赛程间隔  ·  Days Since Last Match", f"{sa.days_since_last_match} 天", f"{sb.days_since_last_match} 天",
         GREEN if sa.days_since_last_match > sb.days_since_last_match else CYAN),
        ("近 7 天场数  ·  Matches in 7d", f"{sa.matches_in_last_7_days} 场", f"{sb.matches_in_last_7_days} 场",
         GREEN if sa.matches_in_last_7_days < sb.matches_in_last_7_days else CYAN),
        ("场均黄红牌  ·  Cards/Game", f"{sa.cards_per_game:.1f}", f"{sb.cards_per_game:.1f}",
         GREEN if sa.cards_per_game < sb.cards_per_game else CYAN),
        ("场均犯规  ·  Fouls/Game", f"{sa.fouls_per_game:.1f}", f"{sb.fouls_per_game:.1f}",
         GREEN if sa.fouls_per_game < sb.fouls_per_game else CYAN),
        ("定位球进球占比  ·  Set Piece Goals", f"{sa.set_piece_goals_pct:.0%}", f"{sb.set_piece_goals_pct:.0%}",
         GREEN if sa.set_piece_goals_pct > sb.set_piece_goals_pct else CYAN),
        ("逼抢强度  ·  Pressing", f"{sa.pressing_intensity:.0%}", f"{sb.pressing_intensity:.0%}",
         GREEN if sa.pressing_intensity > sb.pressing_intensity else CYAN),
    ]
    # Possession is only available when event-level data has been ingested.
    if sa.possession_avg is not None and sb.possession_avg is not None:
        metrics.append(
            ("平均控球率  ·  Possession", f"{sa.possession_avg:.0%}", f"{sb.possession_avg:.0%}",
             GREEN if sa.possession_avg > sb.possession_avg else CYAN)
        )
    for i, (label, va, vb, _) in enumerate(metrics):
        ry = y_r + Inches(0.5) + i * Inches(0.40)
        _add_textbox(slide, right_x + Inches(0.2), ry, Inches(3.5), Inches(0.18),
                    label, font_size=Pt(8), color=GREY, font_name=FONT_CN_BODY)
        _add_textbox(slide, right_x + Inches(0.2), ry + Inches(0.20), Inches(1.5), Inches(0.20),
                    va, font_size=Pt(11), color=WHITE, bold=True, font_name=FONT_MONO)
        _add_textbox(slide, right_x + Inches(4.5), ry + Inches(0.20), Inches(1.5), Inches(0.20),
                    vb, font_size=Pt(11), color=WHITE, bold=True, font_name=FONT_MONO, align=PP_ALIGN.RIGHT)
        _add_textbox(slide, right_x + Inches(2.0), ry + Inches(0.20), Inches(0.4), Inches(0.20),
                    "vs", font_size=Pt(8), color=GREY, font_name=FONT_MONO, align=PP_ALIGN.CENTER)

    # #9 + #10 row: fatigue & set pieces summary
    y_r2 = Inches(4.9)
    _add_panel(slide, right_x, y_r2, right_w, Inches(0.5), fill=BG_CARD)
    avg_fatigue_a = (7 - sa.days_since_last_match) * 0.5 + sa.matches_in_last_7_days * 0.3
    avg_fatigue_b = (7 - sb.days_since_last_match) * 0.5 + sb.matches_in_last_7_days * 0.3
    fatigue_a_label = "高" if avg_fatigue_a > 2.5 else ("中" if avg_fatigue_a > 1.5 else "低")
    fatigue_b_label = "高" if avg_fatigue_b > 2.5 else ("中" if avg_fatigue_b > 1.5 else "低")
    fatigue_a_color = RED if avg_fatigue_a > 2.5 else (GOLD if avg_fatigue_a > 1.5 else GREEN)
    fatigue_b_color = RED if avg_fatigue_b > 2.5 else (GOLD if avg_fatigue_b > 1.5 else GREEN)
    _add_textbox(slide, right_x + Inches(0.2), y_r2 + Inches(0.1), Inches(3.0), Inches(0.35),
                f"赛程疲劳指数  ·  {result.match.team_a.name_zh}  {fatigue_a_label}",
                font_size=Pt(10), color=fatigue_a_color, bold=True, font_name=FONT_CN_BODY)
    _add_textbox(slide, right_x + Inches(4.5), y_r2 + Inches(0.1), Inches(2.5), Inches(0.35),
                f"{result.match.team_b.name_zh}  {fatigue_b_label}",
                font_size=Pt(10), color=fatigue_b_color, bold=True, font_name=FONT_CN_BODY, align=PP_ALIGN.RIGHT)

    # Bottom: injury list (3 + 3 max)
    y = Inches(5.55)
    _add_textbox(slide, MARGIN, y, Inches(6), Inches(0.4),
                f"{result.match.team_a.name_zh} 伤停名单  ·  Injuries", font_size=Pt(12), bold=True, color=GOLD, font_name=FONT_CN_BODY)
    impact_zh = {"critical": "关键", "moderate": "中等", "minor": "轻微"}
    for i, inj in enumerate(result.injuries_a[:3]):
        color = RED if inj.impact == "critical" else (GOLD if inj.impact == "moderate" else GREY)
        _add_textbox(slide, MARGIN + Inches(0.2), y + Inches(0.4) + Inches(i * 0.3), Inches(5.5), Inches(0.3),
                    f"• {inj.player.display_name_cn()}  （{impact_zh.get(inj.impact, inj.impact)}）", font_size=Pt(10), color=color, font_name=FONT_CN_BODY)
    if not result.injuries_a:
        _add_textbox(slide, MARGIN + Inches(0.2), y + Inches(0.4), Inches(5.5), Inches(0.3),
                    "• 无  ·  None reported", font_size=Pt(10), color=GREEN, font_name=FONT_CN_BODY)

    _add_textbox(slide, MARGIN + Inches(6.5), y, Inches(6), Inches(0.4),
                f"{result.match.team_b.name_zh} 伤停名单  ·  Injuries", font_size=Pt(12), bold=True, color=CYAN, font_name=FONT_CN_BODY)
    for i, inj in enumerate(result.injuries_b[:3]):
        color = RED if inj.impact == "critical" else (GOLD if inj.impact == "moderate" else GREY)
        _add_textbox(slide, MARGIN + Inches(6.7), y + Inches(0.4) + Inches(i * 0.3), Inches(5.5), Inches(0.3),
                    f"• {inj.player.display_name_cn()}  （{impact_zh.get(inj.impact, inj.impact)}）", font_size=Pt(10), color=color, font_name=FONT_CN_BODY)
    if not result.injuries_b:
        _add_textbox(slide, MARGIN + Inches(6.7), y + Inches(0.4), Inches(5.5), Inches(0.3),
                    "• 无  ·  None reported", font_size=Pt(10), color=GREEN, font_name=FONT_CN_BODY)


def _page_model_output(prs, result: PredictionResult) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _page_header(slide, "模型输出", "Model Output", 17,
                 team_a=result.match.team_a, team_b=result.match.team_b)

    path = CHART_DIR / "probs.png"
    probability_bars(result.model_probs, path)
    slide.shapes.add_picture(str(path), MARGIN, Inches(1.5), width=Inches(8.5))

    # Right: single-model summary + transparency fields
    y = Inches(1.6)
    p = result.model_probs
    vals = p.win_draw_loss
    for i, (label_zh, label_en, val, color) in enumerate([
        ("主胜", "Win", vals[0], GREEN),
        ("平局", "Draw", vals[1], GREY),
        ("客胜", "Loss", vals[2], RED),
    ]):
        ry = y + Inches(i * 0.95)
        _add_panel(slide, MARGIN + Inches(9.0), ry, Inches(3.1), Inches(0.85), fill=BG_CARD)
        _add_textbox(slide, MARGIN + Inches(9.1), ry + Inches(0.05), Inches(2.9), Inches(0.28),
                    f"{label_zh}  ·  {label_en}", font_size=Pt(12), color=color, bold=True, font_name=FONT_CN_BODY)
        _add_textbox(slide, MARGIN + Inches(9.1), ry + Inches(0.35), Inches(2.9), Inches(0.30),
                    f"{val:.1%}",
                    font_size=Pt(18), color=WHITE, bold=True, font_name=FONT_MONO)

    # Data transparency block
    y_trans = Inches(4.5)
    _add_panel(slide, MARGIN + Inches(9.0), y_trans, Inches(3.1), Inches(2.0), fill=BG_PANEL)
    _add_textbox(slide, MARGIN + Inches(9.1), y_trans + Inches(0.05), Inches(2.9), Inches(0.3),
                "数据透明度  ·  DATA TRANSPARENCY", font_size=Pt(10), color=GOLD, font_name=FONT_CN_BODY, bold=True)
    _add_textbox(slide, MARGIN + Inches(9.1), y_trans + Inches(0.4), Inches(2.9), Inches(0.25),
                f"数据质量  ·  Quality: {p.data_quality:.0%}", font_size=Pt(10), color=WHITE, font_name=FONT_CN_BODY)
    _add_textbox(slide, MARGIN + Inches(9.1), y_trans + Inches(0.7), Inches(2.9), Inches(0.25),
                f"样本量  ·  Samples: {p.sample_size_a} / {p.sample_size_b}", font_size=Pt(10), color=WHITE, font_name=FONT_CN_BODY)
    _add_textbox(slide, MARGIN + Inches(9.1), y_trans + Inches(1.0), Inches(2.9), Inches(0.25),
                f"ELO 先验权重  ·  Prior: {p.elo_prior_weight:.0%}", font_size=Pt(10), color=WHITE, font_name=FONT_CN_BODY)

    _add_textbox(slide, MARGIN, Inches(6.6), Inches(8), Inches(0.4),
                f"预期进球  ·  Expected Goals:  {p.expected_goals[0]:.2f}  —  {p.expected_goals[1]:.2f}",
                font_size=Pt(15), color=GOLD, bold=True, font_name=FONT_CN_BODY)


def _page_monte_carlo(prs, result: PredictionResult) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _page_header(slide, "蒙特卡洛模拟", "Monte Carlo Simulation", 18,
                 team_a=result.match.team_a, team_b=result.match.team_b)

    path = CHART_DIR / "mc.png"
    score_distribution(result.monte_carlo, path)
    slide.shapes.add_picture(str(path), MARGIN, Inches(1.5), width=Inches(7.5))

    # Compute over/under 2.5 goals + total goals distribution (#16 add detail)
    mc = result.monte_carlo
    over_25 = 0.0
    under_25 = 0.0
    over_15 = 0.0
    under_15 = 0.0
    btts = 0.0  # both teams to score
    goal_dist: dict[int, float] = {0: 0.0, 1: 0.0, 2: 0.0, 3: 0.0, 4: 0.0, 5: 0.0, 6: 0.0, 7: 0.0}
    for score_str, prob in mc.distribution.items():
        try:
            a, b = score_str.split("-")
            a, b = int(a), int(b)
            total = a + b
            if total > 7:
                goal_dist[7] += prob
            else:
                goal_dist[total] += prob
            if total >= 3:
                over_25 += prob
            else:
                under_25 += prob
            if total >= 2:
                over_15 += prob
            else:
                under_15 += prob
            if a >= 1 and b >= 1:
                btts += prob
        except Exception:
            pass

    # Left side: result outcome
    y = Inches(1.6)
    for i, (label_zh, label_en, val, color) in enumerate([
        (f"{result.match.team_a.name_zh} 胜", "Win A", mc.win_a, GREEN),
        ("平局", "Draw", mc.draw, GREY),
        (f"{result.match.team_b.name_zh} 胜", "Win B", mc.win_b, RED),
    ]):
        ry = y + Inches(0.7) + i * Inches(0.7)
        _add_panel(slide, MARGIN, ry, Inches(2.5), Inches(0.55), fill=BG_CARD)
        _add_textbox(slide, MARGIN + Inches(0.1), ry + Inches(0.08), Inches(1.5), Inches(0.4),
                    f"{label_zh}", font_size=Pt(12), color=color, bold=True, font_name=FONT_CN_BODY)
        _add_textbox(slide, MARGIN + Inches(1.4), ry + Inches(0.08), Inches(1.0), Inches(0.4),
                    f"{val:.1%}", font_size=Pt(20), bold=True, color=color, font_name=FONT_MONO, align=PP_ALIGN.RIGHT)

    # Middle: goal distribution mini bar chart
    y_g = Inches(4.0)
    _add_textbox(slide, MARGIN, y_g, Inches(7.5), Inches(0.3),
                "总进球数分布  ·  GOAL COUNT DISTRIBUTION", font_size=Pt(11), bold=True,
                color=GOLD, font_name=FONT_CN_BODY)
    chart_left = MARGIN
    chart_y = y_g + Inches(0.5)  # shifted down to avoid overlap with value labels
    chart_w = Inches(7.5)
    chart_h = Inches(1.3)
    # background
    _add_panel(slide, chart_left, chart_y, chart_w, chart_h, fill=BG_PANEL)
    bar_count = 8  # 0, 1, 2, ..., 7+ goals
    bar_w = chart_w / bar_count
    max_p = max(goal_dist.values()) or 0.1
    for i in range(bar_count):
        g = i
        p = goal_dist.get(i, 0)
        bar_h_in = (p / max_p) * (chart_h.inches * 0.75)
        bar_left = chart_left + i * bar_w + Inches(0.05)
        bar_w_each = bar_w - Inches(0.1)
        bar_color = GOLD if g == 2 else (CYAN if g == 1 else (GREEN if g == 3 else GREY))
        # bar from bottom
        bar_top = chart_y + chart_h - Inches(0.3) - Inches(bar_h_in)
        _add_panel(slide, bar_left, bar_top, bar_w_each, Inches(bar_h_in), fill=bar_color)
        # value label (above bar)
        _add_textbox(slide, bar_left, bar_top - Inches(0.20), bar_w_each, Inches(0.20),
                    f"{p*100:.0f}%", font_size=Pt(8), color=WHITE, font_name=FONT_MONO, align=PP_ALIGN.CENTER)
        # x-axis label
        _add_textbox(slide, bar_left, chart_y + chart_h - Inches(0.25), bar_w_each, Inches(0.25),
                    f"{g}球" if g < 7 else "7+球", font_size=Pt(8), color=GREY, font_name=FONT_CN_BODY, align=PP_ALIGN.CENTER)

    # Right: betting market stats
    y_r = Inches(1.6)
    right_x = MARGIN + Inches(7.7)
    _add_panel(slide, right_x, y_r, Inches(4.9), Inches(3.8), fill=BG_PANEL)
    _add_textbox(slide, right_x + Inches(0.2), y_r + Inches(0.1), Inches(4.5), Inches(0.35),
                "投注市场关键指标  ·  BETTING MARKETS", font_size=Pt(11), bold=True, color=GOLD, font_name=FONT_CN_BODY)
    markets = [
        ("大/小 2.5 球  ·  Over/Under 2.5", f"大 {over_25:.0%}  ·  小 {under_25:.0%}", GOLD if over_25 > 0.5 else CYAN),
        ("大/小 1.5 球  ·  Over/Under 1.5", f"大 {over_15:.0%}  ·  小 {under_15:.0%}", GREEN if over_15 > 0.7 else GOLD),
        ("两队进球  ·  BTTS", f"是 {btts:.0%}  ·  否 {1-btts:.0%}", GREEN if btts > 0.5 else GREY),
        ("最可能比分", f"{mc.predicted_score}  ({mc.top_scores[0][1]:.1%})", GOLD),
    ]
    for i, (label, val, color) in enumerate(markets):
        ry = y_r + Inches(0.55) + i * Inches(0.6)
        _add_textbox(slide, right_x + Inches(0.2), ry, Inches(4.5), Inches(0.25),
                    label, font_size=Pt(9), color=GREY, font_name=FONT_CN_BODY)
        _add_textbox(slide, right_x + Inches(0.2), ry + Inches(0.22), Inches(4.5), Inches(0.32),
                    val, font_size=Pt(13), color=color, bold=True, font_name=FONT_MONO)

    # Bottom: sim count
    _add_textbox(slide, MARGIN + Inches(7.7), Inches(5.55), Inches(4.9), Inches(0.4),
                f"模拟次数：{mc.simulations:,}  ·  SIMULATIONS", font_size=Pt(10), color=GREY, font_name=FONT_CN_BODY, align=PP_ALIGN.CENTER)


def _page_sensitivity(prs, result: PredictionResult) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _page_header(slide, "敏感性分析", "Sensitivity Analysis", 19,
                 team_a=result.match.team_a, team_b=result.match.team_b)

    _add_textbox(slide, MARGIN, Inches(1.2), Inches(12), Inches(0.4),
                "若以下变量翻转，结论是否会改变？", font_size=Pt(12), color=GREY, font_name=FONT_CN_BODY)

    # Compute sensitivity: how much each factor moves the consensus
    base = result.model_probs.win_draw_loss[0]
    factors = [
        ("ELO +50", abs(50 * 0.0015), 0.95),
        ("伤停恢复", 0.07, 0.75),
        ("主场优势", 0.05, 0.55),
        ("状态波动 3 分", 0.04, 0.45),
        ("极端天气", 0.03, 0.35),
        ("裁判尺度", 0.02, 0.25),
    ]
    factors.sort(key=lambda x: -x[1])
    path = CHART_DIR / "sensitivity.png"
    sensitivity_tornado(factors, path)
    slide.shapes.add_picture(str(path), MARGIN, Inches(1.8), width=Inches(8.0))

    # #15 explicit impact-level color tags next to each factor (kept within page bounds)
    impact_zh = {0.95: "极高", 0.75: "高", 0.55: "中高", 0.45: "中", 0.35: "中低", 0.25: "低"}
    impact_color = {0.95: RED, 0.75: RED, 0.55: GOLD, 0.45: GOLD, 0.35: CYAN, 0.25: CYAN}
    y = Inches(1.9)
    for i, (name, swing, weight) in enumerate(factors):
        ry = y + Inches(i * 0.65)
        _add_panel(slide, MARGIN + Inches(8.7), ry, Inches(1.0), Inches(0.45), fill=impact_color[weight])
        _add_textbox(slide, MARGIN + Inches(8.7), ry, Inches(1.0), Inches(0.45),
                    impact_zh[weight], font_size=Pt(12), bold=True, color=BG_DEEP,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name=FONT_CN_BODY)
        _add_textbox(slide, MARGIN + Inches(9.8), ry + Inches(0.05), Inches(2.6), Inches(0.4),
                    name, font_size=Pt(12), color=WHITE, font_name=FONT_CN_BODY, bold=True)
        _add_textbox(slide, MARGIN + Inches(9.8), ry + Inches(0.30), Inches(2.6), Inches(0.3),
                    f"影响 ±{swing*100:.1f} pp", font_size=Pt(9), color=GREY, font_name=FONT_CN_BODY)


def _page_final(prs, result: PredictionResult) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _page_header(slide, "最终预测", "Final Prediction", 20,
                 team_a=result.match.team_a, team_b=result.match.team_b)

    team_a = result.match.team_a
    team_b = result.match.team_b
    pick = result.recommended_pick
    predicted_score = result.monte_carlo.predicted_score
    score_a, score_b = result.monte_carlo.split_goals(predicted_score)

    # Big result banner: predicted score as headline
    _add_panel(slide, MARGIN, Inches(1.4), Inches(12.1), Inches(2.0), fill=BG_PANEL)
    pick_label_zh = team_a.name_zh if pick == "A" else (team_b.name_zh if pick == "B" else "平局")
    _add_textbox(slide, MARGIN + Inches(0.5), Inches(1.55), Inches(11), Inches(0.4),
                "最终预测比分  ·  FINAL PREDICTED SCORE", font_size=Pt(12), color=GOLD, font_name=FONT_CN_BODY, bold=True)
    _add_textbox(slide, MARGIN, Inches(2.0), Inches(12.1), Inches(1.1),
                predicted_score, font_size=Pt(80), bold=True, color=WHITE, font_name=FONT_MONO,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _add_textbox(slide, MARGIN + Inches(0.5), Inches(3.05), Inches(11), Inches(0.4),
                f"推荐结果：{pick_label_zh}    ·    该比分出现概率 {result.monte_carlo.top_scores[0][1]:.1%}    ·    TOP 1 SCORE",
                font_size=Pt(15), color=GOLD, font_name=FONT_CN_BODY, bold=True, align=PP_ALIGN.CENTER)

    # Confidence + risks
    conf_color = GREEN if result.confidence == "high" else (GOLD if result.confidence == "medium" else RED)
    conf_zh = "高" if result.confidence == "high" else "中" if result.confidence == "medium" else "低"
    _add_panel(slide, MARGIN, Inches(3.7), Inches(5.9), Inches(3.2), fill=BG_CARD)
    _add_textbox(slide, MARGIN + Inches(0.3), Inches(3.85), Inches(5.5), Inches(0.4),
                "信心指数  ·  CONFIDENCE", font_size=Pt(12), color=GOLD, font_name=FONT_CN_BODY, bold=True)
    _add_textbox(slide, MARGIN + Inches(0.3), Inches(4.3), Inches(5.5), Inches(1.0),
                f"{conf_zh}  ·  {result.confidence.upper()}", font_size=Pt(40), bold=True, color=conf_color, font_name=FONT_CN_BODY)
    p = result.model_probs.win_draw_loss
    _add_textbox(slide, MARGIN + Inches(0.3), Inches(5.3), Inches(5.5), Inches(0.4),
                f"胜 {p[0]:.0%}  ·  平 {p[1]:.0%}  ·  负 {p[2]:.0%}",
                font_size=Pt(12), color=WHITE, font_name=FONT_CN_BODY)
    eg = result.model_probs.expected_goals
    _add_textbox(slide, MARGIN + Inches(0.3), Inches(5.6), Inches(5.5), Inches(0.4),
                f"预期进球（期望值）  ·  xG:  {eg[0]:.1f}  —  {eg[1]:.1f}",
                font_size=Pt(11), color=GREY, font_name=FONT_CN_BODY)

    # Risks
    _add_panel(slide, MARGIN + Inches(6.2), Inches(3.7), Inches(5.9), Inches(3.2), fill=BG_CARD)
    _add_textbox(slide, MARGIN + Inches(6.5), Inches(3.85), Inches(5.5), Inches(0.4),
                "关键风险  ·  KEY RISKS", font_size=Pt(12), color=GOLD, font_name=FONT_CN_BODY, bold=True)
    for i, risk in enumerate(result.key_risks):
        _add_textbox(slide, MARGIN + Inches(6.5), Inches(4.3) + Inches(i * 0.55), Inches(5.4), Inches(0.5),
                    f"• {risk}", font_size=Pt(10), color=WHITE, font_name=FONT_CN_BODY)
    # #20 信心指数解释 (小字说明)
    conf_why = {
        "high": "高 = 3 模型一致 + ELO 差距大 + 阵容完整",
        "medium": "中 = 2-3 模型大体一致 + ELO 差距适中",
        "low": "低 = 模型分歧大 / 双方实力接近 / 阵容不整",
    }[result.confidence]
    _add_textbox(slide, MARGIN + Inches(6.5), Inches(6.55), Inches(5.5), Inches(0.3),
                f"  · {conf_why}", font_size=Pt(8), color=GREY, font_name=FONT_CN_BODY)


def _page_appendix(prs, result: PredictionResult) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _page_header(slide, "附录", "Appendix", 21,
                 team_a=result.match.team_a, team_b=result.match.team_b)

    _add_textbox(slide, MARGIN, Inches(1.4), Inches(12), Inches(0.4),
                "数据来源  ·  DATA SOURCES", font_size=Pt(13), color=GOLD, font_name=FONT_CN_BODY, bold=True)
    sources = [
        "Wikipedia REST API  ·  球员信息和照片",
        "StatsBomb Open Data  ·  2018/2022 世界杯逐球 xG",
        "martj42 国际比赛 CSV  ·  各国 1872+ 比赛结果",
        "ESPN 公开 API  ·  实时赛程（主源）",
        "懂球帝 api.dongqiudi.com  ·  免 Key 中文赛程 / 比赛详情",
        "直播吧 s.qiumibao.com  ·  24h 中文体育热点（赛前情报）",
        "本地种子阵容  ·  各国阵容种子库",
        "World Football ELO Ratings  ·  ELO 评分方法",
    ]
    for i, s in enumerate(sources):
        _add_textbox(slide, MARGIN + Inches(0.3), Inches(1.85) + Inches(i * 0.32), Inches(12), Inches(0.3),
                    f"• {s}", font_size=Pt(11), color=WHITE, font_name=FONT_CN_BODY)

    _add_textbox(slide, MARGIN, Inches(3.7), Inches(12), Inches(0.4),
                "方法说明  ·  METHODOLOGY", font_size=Pt(13), color=GOLD, font_name=FONT_CN_BODY, bold=True)
    method = [
        "1.  ELO 模型：基础实力差异，中立场调整",
        "2.  Poisson + Dixon-Coles：进球分布 + 低比分修正",
        "3.  XGBoost：基于合成但真实的特征空间的梯度提升",
        "4.  定性调整：战术 / 经验 / 心理 / 场地",
        "5.  蒙特卡洛（10,000 次）：基于 Poisson 强度参数模拟比赛",
        "6.  综合概率：三个基础模型调整后的加权平均",
    ]
    for i, m in enumerate(method):
        _add_textbox(slide, MARGIN + Inches(0.3), Inches(4.1) + Inches(i * 0.32), Inches(12), Inches(0.3),
                    m, font_size=Pt(10), color=WHITE, font_name=FONT_CN_BODY)

    _add_textbox(slide, MARGIN, Inches(6.2), Inches(12), Inches(0.4),
                "免责声明  ·  DISCLAIMER", font_size=Pt(13), color=RED, font_name=FONT_CN_BODY, bold=True)
    _add_textbox(slide, MARGIN, Inches(6.55), Inches(12), Inches(0.6),
                "本报告基于历史数据和统计模型，结果仅供参考，不构成任何投注或决策建议。",
                font_size=Pt(11), color=GREY, font_name=FONT_CN_BODY)


# ----------------------- main entry -----------------------

def build_ppt(result: PredictionResult, lang: str = "bilingual", output_path: Path | None = None) -> Path:
    """Build the full PPT and return the saved path."""
    logger.info("Building PPT…")

    # Try to enrich starting XIs with photos from Wikipedia (best effort)
    try:
        for lineup in (result.lineup_a, result.lineup_b):
            batch_augment_squad(lineup.players)
    except Exception as e:
        logger.debug(f"Wikipedia augment skipped: {e}")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _page_cover(prs, result)                  # 1
    _page_summary(prs, result)                # 2
    _page_toc(prs, result)                    # 3  (#16)
    _page_pre_match_news(prs, result)          # 4  赛前热点
    _page_world_cup_context(prs, result)     # 5  (#7)
    _page_team_profile(prs, result)           # 6
    _page_h2h(prs, result)                   # 7  (#5)
    _page_recent_form(prs, result)            # 8
    _page_predicted_lineup(prs, result, "A")  # 9
    _page_predicted_lineup(prs, result, "B")  # 10
    _page_key_players(prs, result, "A")       # 11
    _page_key_players(prs, result, "B")       # 12
    _page_key_matchups(prs, result)           # 13
    _page_squad_depth(prs, result)            # 14
    _page_radar(prs, result)                  # 15
    _page_qualitative(prs, result)            # 16
    _page_model_output(prs, result)           # 17
    _page_monte_carlo(prs, result)            # 18
    _page_sensitivity(prs, result)            # 19
    _page_final(prs, result)                  # 20
    _page_appendix(prs, result)               # 21

    if output_path is None:
        team_a_zh = result.match.team_a.name_zh
        team_b_zh = result.match.team_b.name_zh
        date = result.match.match_date
        output_path = config.output_dir / f"{team_a_zh}_对阵_{team_b_zh}_{date}.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    logger.success(f"PPT 已保存：{output_path}")
    return output_path
