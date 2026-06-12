"""Deeper content audit: text quality, redundancy, gaps."""
from pptx import Presentation
from collections import Counter
import re

prs = Presentation("输出/加拿大_对阵_波黑_2026-06-13.pptx")

# 1. Look for repeated/duplicate text across pages (redundancy)
all_text_by_page = []
for i, slide in enumerate(prs.slides, 1):
    page_texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                page_texts.append(t)
    all_text_by_page.append((i, page_texts))

# Find phrases that appear 3+ times across pages
phrase_counter = Counter()
for _, texts in all_text_by_page:
    for t in texts:
        # Extract meaningful chunks
        for phrase in re.findall(r'[一-鿿]{2,8}', t):
            if len(phrase) >= 2:
                phrase_counter[phrase] += 1

# 2. Find pages with very few shapes (might be sparse)
print("=" * 60)
print("Page shape counts (sortable by content density):")
print("=" * 60)
for i, texts in all_text_by_page:
    n_shapes = len(prs.slides[i-1].shapes)
    n_chars = sum(len(t) for t in texts)
    print(f"  P{i:2d}: {n_shapes} shapes, {n_chars} chars")

# 3. Find pages with "TBD" or placeholder text
print("\n" + "=" * 60)
print("Pages with 'TBD' or placeholder:")
print("=" * 60)
for i, texts in all_text_by_page:
    for t in texts:
        if 'TBD' in t or 'TODO' in t or 'PLACEHOLDER' in t.upper() or '未知' in t:
            print(f"  P{i}: {t[:80]}")
            break

# 4. Phrase repetition (might indicate overuse)
print("\n" + "=" * 60)
print("Top repeated phrases (≥4 occurrences across PPT):")
print("=" * 60)
for phrase, count in phrase_counter.most_common(15):
    if count >= 4:
        print(f"  '{phrase}' x{count}")
